from pptx import Presentation
from pptx.util import Inches, Pt

# Creating Object
ppt = Presentation()

# To create blank slide layout
# We have to use 6 as an argument
# of slide_layouts
blank_slide_layout = ppt.slide_layouts[6]

# Attaching slide obj to slide
slide = ppt.slides.add_slide(blank_slide_layout)

# For adjusting the Margins in inches
left = top = width = height = Inches(1)

# creating textBox
txBox = slide.shapes.add_textbox(left, top,
								width, height)

# creating textFrames
tf = txBox.text_frame
tf.text = "This is text inside a textbox"

# adding Paragraphs
p = tf.add_paragraph()

# adding text
p.text = "This is a second paragraph that's bold and italic"

# font
p.font.bold = True
p.font.italic = True

p = tf.add_paragraph()
p.text = "This is a third paragraph that's big "
p.font.size = Pt(40)

# save file
ppt.save('test_2.pptx')

print("done")
